"""
AquaClear — Extended 30-Slide Presentation Generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (matching the web app Dark Theme) ---
BG_DARK = RGBColor(10, 14, 26)       
BG_CARD = RGBColor(26, 31, 53)       
TEXT_LIGHT = RGBColor(240, 244, 255) 
TEXT_SEC = RGBColor(148, 163, 184)   
ACCENT_OCEAN = RGBColor(6, 182, 212) 
ACCENT_BLUE = RGBColor(59, 130, 246) 
ACCENT_RED = RGBColor(239, 68, 68)   

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_counter = [0]

def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    slide_counter[0] += 1
    return slide

def add_text_box(slide, left, top, width, height, text, font_size=20, color=TEXT_LIGHT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.alignment = align
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def slide_header(slide, title, subtitle=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.1), Inches(0.5), ACCENT_OCEAN)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7), title, font_size=32, color=TEXT_LIGHT, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.5), subtitle, font_size=16, color=TEXT_SEC)
    # Slide number
    add_text_box(slide, prs.slide_width - Inches(1), prs.slide_height - Inches(0.5), Inches(1), Inches(0.5), f"{slide_counter[0]}", font_size=12, color=TEXT_SEC, align=PP_ALIGN.RIGHT)

def add_card(slide, left, top, width, height, title, content, title_color=ACCENT_OCEAN, line_color=None):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, BG_CARD, line_color=line_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5), title, font_size=22, bold=True, color=title_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), height - Inches(1), content, font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 1. TITLE SLIDE
# -----------------------------------------------------------------------------
slide = new_slide()
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), prs.slide_width, Inches(0.05), ACCENT_OCEAN)
add_text_box(slide, Inches(1), Inches(2), Inches(11.33), Inches(1.5), "AquaClear", font_size=70, color=ACCENT_OCEAN, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.7), Inches(11.33), Inches(1), "Advanced Underwater Image Restoration System", font_size=28, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11.33), Inches(0.5), "Khaled Noaman", font_size=20, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.0), Inches(11.33), Inches(0.5), "CENG 455 Digital Image Processing • Istanbul Arel University", font_size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
# 2. INTRODUCTION & MOTIVATION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Introduction & Motivation", "The necessity of underwater computer vision.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Why It Matters", 
         "Ocean exploration heavily relies on optical data for:\n\n"
         "• Marine biology surveying and ecological tracking.\n"
         "• Underwater pipeline and infrastructure inspection.\n"
         "• Autonomous Underwater Vehicle (AUV) navigation.\n\n"
         "Without clear visual data, automated systems fail.")
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "The Optical Challenge", 
         "Water is a dense, hostile medium for optics.\n\n"
         "Light behaves drastically differently than in the atmosphere, leading to catastrophic visual degradation within merely 5 to 10 meters of depth. Standard terrestrial cameras produce unusable data in deep sea scenarios without severe mathematical post-processing.")

# -----------------------------------------------------------------------------
# 3. PROBLEM: COLOR ATTENUATION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "The Problem: Wavelength Attenuation", "Physics of light absorption in water.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Beer-Lambert's Law", 
         "Absorption converts photonic energy into thermal energy. Water absorbs specific frequencies exponentially according to distance.\n\n"
         "• Red Light (650nm) operates with the highest attenuation coefficient. It disappears almost entirely by 5 meters depth.\n"
         "• Orange and Yellow disappear by 15-20 meters.\n"
         "• Blue and Green (short wavelengths) penetrate the deepest.\n\n"
         "This physical law is the sole reason uncorrected underwater photographs possess an overwhelming monochromatic cyan or green color cast.", title_color=ACCENT_RED)

# -----------------------------------------------------------------------------
# 4. PROBLEM: SCATTERING
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "The Problem: Scattering", "The Jaffe-McGlamery Optical Model.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Forward Scattering (Blur)", 
         "Light bouncing off an object is deflected slightly by suspended particles before hitting the camera lens.\n\n"
         "This causes the visual rays to 'spread out'.\n"
         "Result: Severe loss of high-frequency details, blurring of sharp edges, and loss of structural resolution.", title_color=TEXT_LIGHT)
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "Backscattering (Haze/Fog)", 
         "Ambient light from the surface hits particles in the water and deflects directly into the camera lens, never hitting the target object.\n\n"
         "Result: A dense, milky 'veil' or fog that severely cripples global contrast and dynamic range.", title_color=TEXT_LIGHT)

# -----------------------------------------------------------------------------
# 5. SOLUTION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "The Solution: Classical Image Processing", "Why not Neural Networks?")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Deterministic Mathematical Modeling", 
         "While deep learning (CNNs/GANs) is popular, AquaClear utilizes purely classical IP matrices. Why?\n\n"
         "1. Data Independence: Deep learning requires massive pairs of dry/water images which are impossible to capture realistically. Classical math requires no training data.\n"
         "2. No Hallucinations: Neural nets often invent or hallucinate textures. Classical algorithms strictly preserve existing physical data, making them safe for scientific metrics.\n"
         "3. Explainability: Every pixel shift in this pipeline can be mathematically proven via physical optical theorems.\n"
         "4. Extreme Efficiency: Processes 4K images in < 0.5 seconds on a CPU.")

# -----------------------------------------------------------------------------
# 6. ARCHITECTURE OVERVIEW
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "System Architecture", "Full-Stack Web Application.")
add_card(slide, Inches(1), Inches(2), Inches(3.5), Inches(4.5), "Client (Frontend)", 
         "Vanilla HTML & CSS\n\n"
         "• Flat UI Design\n"
         "• Dark/Light Mode Toggle\n"
         "• Asynchronous JS API calls\n"
         "• Zero heavy frameworks", title_color=ACCENT_OCEAN)
add_card(slide, Inches(4.9), Inches(2), Inches(3.5), Inches(4.5), "Router (Backend)", 
         "Node.js & Express\n\n"
         "• Handles HTTP stateless requests\n"
         "• Multer for multipart form uploads\n"
         "• Forks Python subprocesses securely", title_color=ACCENT_BLUE)
add_card(slide, Inches(8.8), Inches(2), Inches(3.5), Inches(4.5), "IP Engine (Core)", 
         "Python 3.10 \n\n"
         "• OpenCV (cv2) for spatial filtering\n"
         "• NumPy for vectorized C-level math calculations\n"
         "• Sub-second matrix processing", title_color=TEXT_LIGHT)

# -----------------------------------------------------------------------------
# 7. THE PIPELINE
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "The Mathematical Pipeline", "Step-by-step sequential processing.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4), "Sequential Flow", 
         "1. Auto-Analysis Heuristics (Evaluate image globally)\n"
         "2. Edge-Preserving Noise Reduction (Bilateral Filter)\n"
         "3. Color Cast Compensation (Gray World Theorem)\n"
         "4. Backscatter Fog Removal (Dark Channel Prior)\n"
         "5. Adaptive Contrast Stretching (CLAHE)\n"
         "6. High-Frequency Sharpening (Unsharp Masking)\n"
         "7. Objective Metric Calculation (PSNR, SSIM, Entropy)")

# -----------------------------------------------------------------------------
# 8. PHASE 1: AUTO-ANALYSIS
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Phase 1: Auto-Analysis Matrix", "Dynamic Parameterization.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Replacing Magic Numbers", 
         "Static pipelines fail because ocean conditions vary wildly. AquaClear dynamically scans the matrix first:\n\n"
         "• Blue-Red Mean Differential: Measures depth. If Red deficit > 30, triggers extreme color correction.\n"
         "• Dark Channel Mean: Measures Haze. If > 0.4, identifies heavy backscatter and increases Dehazing Omega.\n"
         "• Laplacian Variance: Measures blur. If < 100, identifies severe forward scattering and switches to Non-Local Means denoising + heavy Sharpening.")

# -----------------------------------------------------------------------------
# 9. AUTO-ANALYSIS CODE
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Auto-Analysis Metrics Logic", "Global diagnostic calculation.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Python / NumPy Heuristics", 
         "b_mean, g_mean, r_mean = cv2.split(img)\n"
         "red_deficit = b_mean.mean() - r_mean.mean()\n\n"
         "min_channel = np.min(img, axis=2).astype(np.float64)\n"
         "dark_mean = min_channel.mean() / 255.0\n\n"
         "gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)\n"
         "laplacian_var = cv2.Laplacian(gray, cv2.CV_64F).var()\n"
         "contrast = gray.std()", title_color=ACCENT_BLUE, line_color=ACCENT_BLUE)

# -----------------------------------------------------------------------------
# 10. PHASE 2: NOISE REDUCTION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Phase 2: Noise Attenuation", "Dealing with Marine Snow.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Linear Filter Failure", 
         "Standard Gaussian or Mean blur filters reduce noise but destroy existing edges.\n\n"
         "Since underwater images are already blurred by forward scattering, using a linear Gaussian blur renders the image permanently destroyed.", title_color=ACCENT_RED)
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "Bilateral Filtering", 
         "A non-linear, edge-preserving spatial filter.\n\n"
         "Combines a domain kernel (spatial distance) with a range kernel (intensity difference).\n"
         "It smooths flat water backgrounds perfectly, but forcefully stops smoothing when it detects a sharp color transition (a fish edge).", title_color=ACCENT_OCEAN)

# -----------------------------------------------------------------------------
# 11. PHASE 3: COLOR CORRECTION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Phase 3: Color Cast Compensation", "Restoring lost wavelengths.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Von Kries Gray World Hypothesis", 
         "Theory states that in a scene with vast variety, the global average of Red, Green, and Blue pixels should equate to neutral gray (R=G=B).\n\n"
         "We calculate the global mean of the entire image, then calculate specific gains for each color channel relative to that mean.\n\n"
         "By multiplying the isolated matrices by their respective gains, the highly attenuated Red channel is multiplied exponentially, balancing the overall color histogram and stripping the cyan tint.")

# -----------------------------------------------------------------------------
# 12. PHASE 4: DEHAZING
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Phase 4: Dark Channel Prior (DCP)", "Subtracting backscattered fog.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Statistical Observation", 
         "He et al. discovered that in clear, outdoor images, non-sky patches always contain pixels with an intensity near zero in at least one color channel (due to shadows or dark colors).\n\n"
         "In underwater images, these 'dark channels' are instead filled with bright intensity. That extra brightness is exclusively the backscattered haze!\n\n"
         "By isolating the dark channel, we objectively isolate the mathematical volume of the physical fog.")

# -----------------------------------------------------------------------------
# 13. DCP TRANSMISSION MAPPING
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Transmission Mapping", "Recovering Radiance.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "The Math", 
         "1. Atmospheric Light (A): Estimated by finding the top 0.1% brightest pixels in the Dark Channel matrix.\n"
         "2. Transmission Map (t): The percentage of light that survived scattering.\n"
         "   t(x) = 1 - ω * [ min(I(y) / A) ]\n"
         "3. Radiance Recovery (J): Subtracting the haze via the degradation inverse function:\n"
         "   J(x) = (I(x) - A) / max(t(x), 0.1) + A", line_color=ACCENT_BLUE)

# -----------------------------------------------------------------------------
# 14. PHASE 5: CONTRAST (CLAHE)
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Phase 5: Adaptive Contrast", "Restoring dynamic range without blowing out.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Global vs Local HE", 
         "Global Histogram Equalization stretches pixel intensities across the full 0-255 spectrum over the entire image. This violently 'washes out' underwater images.\n\n"
         "Adaptive HE divides the image into 8x8 grids, equalizing textures locally.", title_color=TEXT_LIGHT)
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "The 'Clip Limit'", 
         "CLAHE (Contrast Limited) caps the histogram peaks before equalizing.\n\n"
         "Without the clip limit, flat areas like open ocean would have tiny sensor noise amplified into massive, ugly blocks of grain.", title_color=ACCENT_OCEAN)

# -----------------------------------------------------------------------------
# 15. LAB COLOR SPACE (UZEM)
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "UZEM Forum Requirement", "Why apply CLAHE to LAB L-Channel?")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Preventing Color Distortion", 
         "Applying equalization directly to R, G, and B matrices independently alters their mathematical ratios, which artificially changes the actual color tracking of the image (making fish look neon, for example).\n\n"
         "By cleanly converting from RGB to LAB (CIE1976), we isolate Luminance (L) from chrominance (A, B). Applying CLAHE exclusively to the L channel stretches the contrast strictly based on brightness, perfectly preserving the restored hues from Phase 3.")

# -----------------------------------------------------------------------------
# 16. PHASE 6: SHARPENING
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Phase 6: Frequency Domain Blending", "Defeating Forward Scattering.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Unsharp Masking", 
         "A high-pass frequency filter used to accentuate edges.\n\n"
         "1. Blur the image (creates a Low-Pass image representing flat colors).\n"
         "2. Subtract the Low-Pass from the Original. The mathematical remainder is strictly the high-frequency edge data (the 'Mask').\n"
         "3. Multiply the Mask by a strength factor (dynamically assigned by Phase 1 Auto-Analysis) and add it back to the original image.\n\n"
         "Edges 'pop', restoring structural clarity.")

# -----------------------------------------------------------------------------
# 17. METRICS OVERVIEW
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Quality Assessment Metrics", "Replacing subjective sight with math.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Objective Verification", 
         "Because human eyes are easily tricked by oversaturated colors, we output 4 mathematical metrics on every processed image:\n\n"
         "1. PSNR (Reconstruction Fidelity)\n"
         "2. SSIM (Perceptual Geometry)\n"
         "3. Colorfulness Index (Saturation Validation)\n"
         "4. Shannon Entropy (Information Distribution Density)")

# -----------------------------------------------------------------------------
# 18. PSNR & SSIM
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Structural Metrics", "PSNR and SSIM.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Peak Signal-to-Noise Ratio", 
         "10 * log10 (MAX² / Mean Squared Error)\n\n"
         "Measures the delta between the raw and enhanced image matrices. High values (20dB - 40dB) prove the algorithms did not mathematically destroy the underlying structural array.", title_color=ACCENT_OCEAN)
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "Structural Similarity (SSIM)", 
         "Evaluates Luminance, Contrast, and Structure independently across sliding windows.\n\n"
         "Ensures that while the color and fog changed, the physical shapes of the subjects did not warp or distort.", title_color=ACCENT_BLUE)

# -----------------------------------------------------------------------------
# 19. ENTROPY
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Mathematical Information", "Shannon Entropy.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Shannon Entropy Calculation", 
         "Entropy = - Σ P(i) * log2 P(i)\n\n"
         "Quantifies the amount of 'information' inside an image. An uncorrected, heavily fogged image has a very narrow band of pixel intensities, giving it low entropy.\n\n"
         "After applying Dehaze and CLAHE algorithms, the 0-255 histogram is massively widened. Shannon Entropy reliably increases across all successfully enhanced outputs, validating the process.")

# -----------------------------------------------------------------------------
# 20. VISUAL RESULTS
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Visual Results Analysis", "Before & After Evaluation.")
try:
    sample_dir = os.path.join(os.path.dirname(__file__), 'sample_images')
    if os.path.exists(os.path.join(sample_dir, 'coral_reef.jpg')):
        slide.shapes.add_picture(os.path.join(sample_dir, 'coral_reef.jpg'), Inches(1), Inches(2), width=Inches(4.5))
        add_text_box(slide, Inches(1), Inches(5), Inches(4.5), Inches(0.5), "Raw Negative Image", font_size=16, color=ACCENT_RED, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(6), Inches(2), Inches(6.33), Inches(4.5), "Shallow Coral Evaluation", 
         "In 2-5m depths, main degradation is color loss.\n\n"
         "The Gray World implementation successfully detects the moderate red deficit and scales it. The algorithm intelligently limits dehazing since backscatter is low, preventing the image from becoming overly darkened. Natural vibrant yellows and oranges are restored.")
except:
    pass

# -----------------------------------------------------------------------------
# 21. DEEP WATER RESULTS
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Extreme Environments", "Deep Ocean Wrecks.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Pelagic & Wreck Evaluation", 
         "In extreme depths, images are almost perfectly monochromatic cyan and completely blocked by backscattered transmission fog.\n\n"
         "The Auto-Diagnostic flags critical red loss and high laplacian blur. The Dark Channel Prior map isolates the massive fog volume perfectly, allowing the inverse function to subtract it. Contrast and visibility ranges are mathematically doubled within 0.5 seconds.")

# -----------------------------------------------------------------------------
# 22. PRESETS AND THE GUI
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "User Interface Controls", "The Web Application Integration.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Manual vs Automatic", 
         "While Auto-Analyze maps the optimal scientific parameters flawlessly, the UI allows human overrides.\n\n"
         "Users can select 'Mild' or 'Strong' structural limits, or enter 'Custom' mode. Custom mode exposes the HTML range sliders attached directly to the Python runtime subprocess parameters. Sliding the 'Dehaze Omega' bar instantly re-computes the physical transmission map matrix.")

# -----------------------------------------------------------------------------
# 23. LIMITATIONS
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "System Limitations", "Edge Cases in Classical Optics.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Gray World Failure", 
         "If the actual physical scene is monochromatic (e.g. an extreme close-up of a giant yellow sponge taking up 95% of the frame), the Gray World math fails catastrophically. It interprets the real yellow as a 'water tint' and violently strips it away, destroying the image.", title_color=ACCENT_RED)
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "Artificial Lighting", 
         "The Dark Channel Prior model assumes ambient atmospheric light is uniform. If a diver introduces a bright, localized flashlight beam, the DCP matrix calculations shatter, resulting in black 'halos' developing around the light source during subtraction.", title_color=ACCENT_RED)

# -----------------------------------------------------------------------------
# 24. TRADEOFFS
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Architectural Tradeoffs", "Balancing Speed vs Quality.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "CLAHE Tile Dimensions", 
         "The primary engineering tradeoff is selecting the block size for the CLAHE algorithm grid.\n\n"
         "Large Blocks: Highly efficient CPU calculation times, but causes severe 'tiling' seam artifacts on the image background.\n"
         "Tiny Blocks: Flawless bilinear interpolation and perfect smoothing, but exponentially multiplies the CPU mathematical workload.\n"
         "Solution: A dynamic midpoint tied to the global diagnostic variance scale.")

# -----------------------------------------------------------------------------
# 25. CONCLUSION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Conclusion", "Final Thoughts.")
add_card(slide, Inches(1), Inches(2), Inches(11.33), Inches(4.5), "Project Summary", 
         "AquaClear successfully engineers a purely classical, math-based image processor that reverses the physics of water degradation.\n\n"
         "By building sequential modules (Bilateral filtering, DCP dehazing, CLAHE), it proves that deterministic algorithms are still vastly superior to black-box Deep Learning for scientific reliability.\n\n"
         "The incorporation of a stateless Node.js server and an adaptive heuristic engine creates a robust, fast, and highly reliable modern application.")

# -----------------------------------------------------------------------------
# 26. FUTURE WORK
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Future Work", "Expanding the Scope.")
add_card(slide, Inches(1), Inches(2), Inches(5.4), Inches(4.5), "Stereo Depth Mapping", 
         "If integrated with dual-lens stereo cameras, we could extract physical depth arrays (z-index) natively. We could then apply exponential attenuation modeling per-pixel with flawless millimeter accuracy, circumventing the need for statistical estimation entirely.", title_color=ACCENT_OCEAN)
add_card(slide, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), "Live Video Telemetry", 
         "By porting the vectorized Python NumPy matrices strictly into C++ with CUDA parallelization architectures, the pipeline could process live 60fps underwater feeds natively, integrating directly into AUV piloting dashboards.", title_color=ACCENT_OCEAN)

# -----------------------------------------------------------------------------
# 27. Q&A
# -----------------------------------------------------------------------------
slide = new_slide()
add_text_box(slide, Inches(1), Inches(3), Inches(11.33), Inches(2), "THANK YOU\nQuestions & Answers", font_size=50, color=ACCENT_OCEAN, bold=True, align=PP_ALIGN.CENTER)

# Save
def setup_borders():
    pass

out_path = os.path.join(os.path.dirname(__file__), 'AquaClear_Extended_30_Slides.pptx')
prs.save(out_path)
print(f"Massive presentation saved to {out_path}")
