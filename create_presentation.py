"""
AquaClear — Underwater Image Enhancement Presentation Generator
Creates a PowerPoint presentation tailored for a live demo.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (matching the web app Dark Theme) ---
BG_DARK = RGBColor(10, 14, 26)       # --bg-primary
BG_CARD = RGBColor(26, 31, 53)       # --bg-card
TEXT_LIGHT = RGBColor(240, 244, 255) # --text-primary
TEXT_SEC = RGBColor(148, 163, 184)   # --text-secondary
ACCENT_OCEAN = RGBColor(6, 182, 212) # Teal/Cyan Accent
ACCENT_BLUE = RGBColor(59, 130, 246) # Blue Accent
ACCENT_RED = RGBColor(239, 68, 68)   # Red for warnings

prs = Presentation()
# Set 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_counter = [0]


def new_slide():
    """Create a new slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    slide_counter[0] += 1
    return slide


def add_text_box(slide, left, top, width, height, text, font_size=20, color=TEXT_LIGHT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Segoe UI'
    p.alignment = align
    return txBox


def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def slide_header(slide, title, subtitle=None):
    """Standard header for all content slides."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.1), Inches(0.5), ACCENT_OCEAN)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 title, font_size=32, color=TEXT_LIGHT, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=TEXT_SEC)


def add_step_badge(slide, step_num, total_steps, text):
    """Badge in top right to track demo progress."""
    badge_width = Inches(3.0)
    left = prs.slide_width - badge_width - Inches(0.5)
    
    # Background pill
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(0.3), badge_width, Inches(0.5), BG_CARD, line_color=ACCENT_OCEAN)
    
    # Progress text
    add_text_box(slide, left, Inches(0.3), badge_width, Inches(0.5),
                 f"Demo Step {step_num}/{total_steps}: {text}",
                 font_size=12, color=ACCENT_OCEAN, bold=True, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
# 1. TITLE SLIDE
# -----------------------------------------------------------------------------
slide = new_slide()
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), prs.slide_width, Inches(0.05), ACCENT_OCEAN)
add_text_box(slide, Inches(1), Inches(2), Inches(11.33), Inches(1.5),
             "AquaClear", font_size=70, color=ACCENT_OCEAN, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.7), Inches(11.33), Inches(1),
             "Underwater Image Enhancement & Color Correction", font_size=28, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11.33), Inches(0.5),
             "Khaled Noaman", font_size=20, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.0), Inches(11.33), Inches(0.5),
             "Image Processing Course • Istanbul Arel University • 2026", font_size=16, color=TEXT_SEC, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
# 2. PROBLEM STATEMENT
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "The Problem: Underwater Vision", "Why do underwater photos look bad?")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(3.5), Inches(4), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(3.1), Inches(0.5), "Color Attenuation", font_size=24, bold=True, color=ACCENT_RED)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(3.1), Inches(3),
             "• Water absorbs light depending on wavelength.\n\n"
             "• Red disappears first (at ~5m).\n\n"
             "• Results in strong blue/green color cast.", font_size=18, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(2), Inches(3.5), Inches(4), BG_CARD)
add_text_box(slide, Inches(5.1), Inches(2.2), Inches(3.1), Inches(0.5), "Scattering & Haze", font_size=24, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(5.1), Inches(2.8), Inches(3.1), Inches(3),
             "• Particles in water scatter light.\n\n"
             "• Creates a foggy 'haze' effect.\n\n"
             "• Lowers overall visibility and contrast.", font_size=18, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2), Inches(3.5), Inches(4), BG_CARD)
add_text_box(slide, Inches(9.0), Inches(2.2), Inches(3.1), Inches(0.5), "Low Contrast & Blur", font_size=24, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(9.0), Inches(2.8), Inches(3.1), Inches(3),
             "• Loss of high-frequency details.\n\n"
             "• Edges become soft.\n\n"
             "• Floating marine snow adds noise.", font_size=18, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 3. SOLUTION: PURE IP APPROACH
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "The Solution: Classical IP", "No Neural Networks, just Math and Matrices.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(2), BG_CARD)
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.33), Inches(0.5), "Why Classical IP?", font_size=24, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10.33), Inches(1),
             "• Fast and lightweight (runs in milliseconds).\n"
             "• Highly interpretable (we know exactly what every parameter does).\n"
             "• Deterministic results (no hallucinations).", font_size=18, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(11.33), Inches(2), BG_CARD)
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.33), Inches(0.5), "The AquaClear Pipeline", font_size=24, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10.33), Inches(1),
             "1. Noise Reduction → 2. Color Correction → 3. Dehazing → 4. Contrast → 5. Sharpening", font_size=20, color=ACCENT_OCEAN, bold=True, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------------
# 4. PHYSICS OF UNDERWATER LIGHT
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Physics of Underwater Light", "Why do colors disappear?")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(5.4), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(0.5), "Wavelength Absorption", font_size=22, bold=True, color=ACCENT_RED)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(5), Inches(3.5),
             "Water acts as a dense filter. It absorbs different wavelengths at different depths:\n\n"
             "• <b>Red (650nm)</b>: Disappears first (3-5 meters).\n"
             "• <b>Orange/Yellow</b>: Fades by 15-20 meters.\n"
             "• <b>Blue/Green</b>: Penetrates the deepest.\n\n"
             "Result = The classic 'washed-out blue' tint.", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5), "The Jaffe-McGlamery Model", font_size=22, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
             "Total Image = Direct Signal + Forward Scatter + Backscatter\n\n"
             "• <b>Direct Signal</b>: The clear light from the object.\n"
             "• <b>Forward Scatter</b>: Blurs the edges.\n"
             "• <b>Backscatter</b>: Hit particles and bounces back (Haze).\n\n"
             "Our algorithms try to reverse this model.", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 5. DIGITAL IMAGE FUNDAMENTALS (IP Theory)
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Digital Image Fundamentals", "Applied theory in the project.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(5.4), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(0.5), "Color Spaces Used", font_size=22, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(5), Inches(3.5),
             "• <b>RGB</b>: Base matrices. Used for display and Gray World compensation.\n"
             "• <b>LAB (CIE1976)</b>: Separates Luminance (L) from Color (A,B). Used for Contrast Enhancement (CLAHE) to preserve perceived hues.\n"
             "• <b>HSV</b>: Used for calculating colorfulness and saturation metrics.", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5), "Spatial Domain Filtering", font_size=22, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
             "Modifying pixels based on their neighborhoods.\n\n"
             "• <b>Convolution</b>: Using n x n spatial kernels.\n"
             "• <b>Laplacian (2nd derivative)</b>: For edge detection.\n"
             "• <b>Gaussian Filters</b>: For smoothing.", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 6. DEEP DIVE: NOISE REDUCTION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Deep Dive: Noise Reduction", "Removing marine snow and sensor noise.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(2), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(0.5), "The Problem: Non-Linear Noise", font_size=22, bold=True, color=ACCENT_RED)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(1),
             "Standard linear filters (like Mean or Gaussian blur) will reduce noise, but they ruin the image by blurring the edges. Underwater images already suffer from blur due to forward scattering.", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(11.33), Inches(2), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(0.5), "The Solution: Bilateral Filtering", font_size=22, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1),
             "Bilateral filters combine a spatial Gaussian (distance) with a range Gaussian (intensity difference).\n"
             "• Smooths flat regions (removes noise).\n"
             "• Stops smoothing when it detects a sharp intensity jump (keeps edges sharp).", font_size=16, color=TEXT_SEC)


# -----------------------------------------------------------------------------
# 7. DEMO START OVERVIEW
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Live Demo: AquaClear Workflow", "How the application processes images.")

add_shape(slide, MSO_SHAPE.OVAL, Inches(2), Inches(3), Inches(2), Inches(2), BG_CARD, line_color=ACCENT_OCEAN)
add_text_box(slide, Inches(2), Inches(3.7), Inches(2), Inches(0.6), "1\nUpload", font_size=20, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.OVAL, Inches(5.66), Inches(3), Inches(2), Inches(2), BG_CARD, line_color=ACCENT_OCEAN)
add_text_box(slide, Inches(5.66), Inches(3.7), Inches(2), Inches(0.6), "2\nAnalyze", font_size=20, bold=True, color=ACCENT_OCEAN, align=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.OVAL, Inches(9.33), Inches(3), Inches(2), Inches(2), BG_CARD, line_color=ACCENT_OCEAN)
add_text_box(slide, Inches(9.33), Inches(3.7), Inches(2), Inches(0.6), "3\nEnhance", font_size=20, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Draw arrows
add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.8), Inches(1), Inches(0.4), ACCENT_OCEAN)
add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7.9), Inches(3.8), Inches(1), Inches(0.4), ACCENT_OCEAN)

# -----------------------------------------------------------------------------
# 8. DEMO STEP 1: UPLOAD & AUTO-ANALYZE
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Step 1: Auto-Analysis Matrix", "Detecting physical properties from pixels.")
add_step_badge(slide, 1, 4, "Auto-Analyze")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(0.5), "Diagnostic Processing", font_size=22, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(11), Inches(1),
             "Before applying filters, the backend calculates global metrics:\n"
             "• <b>Color Shift:</b> Mean(B) - Mean(R) > Threshold indicates depth.\n"
             "• <b>Contrast:</b> Global standard deviation (if < 30, needs heavy CLAHE).\n"
             "• <b>Blur:</b> Variance of Laplacian (if < 100, needs high unsharp strength).", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(11.33), Inches(2.5), BG_CARD, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(4.2), Inches(11), Inches(0.5), "Code Snippet: Dynamic Parameterization", font_size=20, bold=True, color=TEXT_LIGHT)
code = """laplacian_var = cv2.Laplacian(gray, cv2.CV_64F).var()

if laplacian_var < 100:
    settings["unsharp_strength"] = 2.5   # Very blurry, strong sharpening
    settings["noise_method"] = "nlm"     # Use Non-Local Means (heavy noise)
elif laplacian_var < 300:
    settings["unsharp_strength"] = 1.5   # Standard sharpening"""
add_text_box(slide, Inches(1.2), Inches(4.8), Inches(11), Inches(1.5), code, font_size=14, color=ACCENT_OCEAN)

# -----------------------------------------------------------------------------
# 9. DEMO STEP 2: COLOR CORRECTION & DEHAZE
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Step 2: Color Shift & Transmission", "Fixing attenuation & scattering.")
add_step_badge(slide, 2, 4, "Color & Dehaze")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(5.4), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(0.5), "Gray World Theorem", font_size=22, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(5), Inches(3.5),
             "Hypothesis: Given an image with enough variety, the average value of R, G, and B should be equal (neutral gray).\n\n"
             "Implementation:\n"
             "1. Calculate mean of R, G, B.\n"
             "2. Find global average.\n"
             "3. Scale up deficient channels (Red) by (Avg / Mean_R).\n"
             "Optional: Channel compensation blends Red based on Green.", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2), Inches(5.4), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5), "Dark Channel Prior (DCP)", font_size=22, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(7.1), Inches(2.8), Inches(5), Inches(3.5),
             "Breakthrough algorithm (He et al.).\n\n"
             "1. <b>Dark Channel:</b> Minimum filter across R,G,B locally.\n"
             "2. <b>A-light:</b> Use brightest 0.1% of Dark Channel pixels to estimate Atmospheric Light.\n"
             "3. <b>Transmission Map (t):</b> Calculates scattering distance.\n"
             "4. <b>Recovery:</b> J = (I - A)/max(t, 0.1) + A.", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 10. DEMO STEP 3: CONTRAST & SHARPENING
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Step 3: Edge Details", "Making features distinct.")
add_step_badge(slide, 3, 4, "Details")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(2), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(0.5), "CLAHE vs Global Histogram Equalization", font_size=22, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(1),
             "• <b>Global HE:</b> Spreads all pixel intensities across the full 0-255 range. It washes out scenes and destroys local contrast.\n"
             "• <b>CLAHE:</b> Divides image into 8x8 blocks. Equalizes locally. Clips peaks in histograms to prevent noise blowout. Re-interpolates using bilinear interpolation to remove grid lines.", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(11.33), Inches(2), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(0.5), "Unsharp Masking (Frequency Domain)", font_size=22, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1),
             "High-frequency data = edges. Low-frequency data = flat colors.\n"
             "1. Create Low-Pass image (Gaussian Blur).\n"
             "2. High-Pass (Edges) = Original - Low-Pass.\n"
             "3. Sharpened = Original + (Strength * High-Pass).", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 11. DEMO STEP 4: METRICS & RESULTS
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Step 4: Objective Quality Assessment", "Math measuring quality.")
add_step_badge(slide, 4, 4, "Metrics")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(5.4), Inches(2.1), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(0.5), "PSNR (Peak Signal-to-Noise)", font_size=18, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(5), Inches(1.2),
             "Measures reconstruction fidelity. 10 * log10(MAX² / Mean Squared Error). Values 20-40dB are typical; higher means less structural damage from filters.", font_size=14, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2), Inches(5.4), Inches(2.1), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.5), "SSIM (Structural Similarity)", font_size=18, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(7.1), Inches(2.7), Inches(5), Inches(1.2),
             "Perceptual metric. Measures structural information change rather than absolute errors in pixels. Considers Luminance, Contrast, and Structure independently.", font_size=14, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.4), Inches(5.4), Inches(2.1), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(4.6), Inches(5), Inches(0.5), "UCIQE (Underwater Color Image Quality)", font_size=18, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(1.2), Inches(5.1), Inches(5), Inches(1.2),
             "Linear combination of standard deviation of chroma, avg saturation, and contrast of luminance. Designed specifically to measure underwater un-blurring.", font_size=14, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.4), Inches(5.4), Inches(2.1), BG_CARD)
add_text_box(slide, Inches(7.1), Inches(4.6), Inches(5), Inches(0.5), "Shannon Entropy", font_size=18, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(7.1), Inches(5.1), Inches(5), Inches(1.2),
             "Measures information content/richness via histogram probability distribution. Enhancement should dramatically increase total entropy.", font_size=14, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 12. SYSTEM ARCHITECTURE
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Complete System Architecture", "Production-ready Web Application.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(4.5), BG_CARD)

# Drawing the Architecture flow
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3), Inches(2.5), Inches(1.5), BG_DARK, line_color=TEXT_SEC)
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(2.5), Inches(1), "User Interface\n(HTML/CSS/JS)", font_size=20, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(3.5), Inches(0.8), Inches(0.5), ACCENT_OCEAN)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(3), Inches(2.5), Inches(1.5), BG_DARK, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(5.2), Inches(3.2), Inches(2.5), Inches(1), "Node.js Server\n(Express/Multer)", font_size=20, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(7.9), Inches(3.5), Inches(0.8), Inches(0.5), ACCENT_OCEAN)

add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8.9), Inches(3), Inches(2.5), Inches(1.5), BG_DARK, line_color=ACCENT_OCEAN)
add_text_box(slide, Inches(8.9), Inches(3.2), Inches(2.5), Inches(1), "Python Engine\n(OpenCV/NumPy)", font_size=20, bold=True, color=ACCENT_OCEAN, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5), Inches(10), Inches(1),
             "• Client interacts with stateless API.\n"
             "• Node.js handles HTTP parsing and temporary file storage.\n"
             "• Python is invoked via Subprocesses (JSON payload in/out). Data returned via Base64 Strings.", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 13. VISUAL RESULTS (BEFORE & AFTER)
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Visual Results", "Before & After Comparisons from the App.")

try:
    sample_dir = os.path.join(os.path.dirname(__file__), 'sample_images')
    # Try to load a real image to show the problem vs solution
    if os.path.exists(os.path.join(sample_dir, 'coral_reef.jpg')):
        slide.shapes.add_picture(os.path.join(sample_dir, 'coral_reef.jpg'), Inches(1), Inches(2), width=Inches(3.3))
        add_text_box(slide, Inches(1), Inches(4.5), Inches(3.3), Inches(0.5), "Coral Reef (Raw)", font_size=14, color=TEXT_SEC, align=PP_ALIGN.CENTER)
    
    if os.path.exists(os.path.join(sample_dir, 'jellyfish.jpg')):
        slide.shapes.add_picture(os.path.join(sample_dir, 'jellyfish.jpg'), Inches(5), Inches(2), width=Inches(3.3))
        add_text_box(slide, Inches(5), Inches(4.5), Inches(3.3), Inches(0.5), "Jellyfish (Raw)", font_size=14, color=TEXT_SEC, align=PP_ALIGN.CENTER)
        
    if os.path.exists(os.path.join(sample_dir, 'deep_ocean.jpg')):
        slide.shapes.add_picture(os.path.join(sample_dir, 'deep_ocean.jpg'), Inches(9), Inches(2), width=Inches(3.3))
        add_text_box(slide, Inches(9), Inches(4.5), Inches(3.3), Inches(0.5), "Deep Ocean (Raw)", font_size=14, color=TEXT_SEC, align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11.33), Inches(1),
                 "The AquaClear pipeline automatically detects the dense blue/green tint in these raw inputs, calculates the atmospheric transmission map, and restores the lost reds and contrast in < 0.5 seconds.", font_size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
except Exception as e:
    add_text_box(slide, Inches(1), Inches(3), Inches(11.33), Inches(1), f"(Images omitted - ensure sample_images exist. Error: {e})", font_size=16, color=ACCENT_RED, align=PP_ALIGN.CENTER)


# -----------------------------------------------------------------------------
# 14. TECH STACK
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Architectural Stack", "How the pieces fit together.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(3.5), Inches(4), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(3.1), Inches(0.5), "Frontend UI", font_size=24, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(3.1), Inches(3),
             "• Vanilla HTML/CSS/JS (Flat Design)\n\n"
             "• Data-theme (Dark/Light)\n\n"
             "• LocalStorage for theme persistence\n\n"
             "• No frontend frameworks", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(2), Inches(3.5), Inches(4), BG_CARD)
add_text_box(slide, Inches(5.1), Inches(2.2), Inches(3.1), Inches(0.5), "Web Server", font_size=24, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(5.1), Inches(2.8), Inches(3.1), Inches(3),
             "• Node.js + Express\n\n"
             "• Multer for multipart form uploads\n\n"
             "• Subprocesses to trigger Python via CLI", font_size=16, color=TEXT_SEC)

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2), Inches(3.5), Inches(4), BG_CARD)
add_text_box(slide, Inches(9.0), Inches(2.2), Inches(3.1), Inches(0.5), "IP Engine", font_size=24, bold=True, color=TEXT_LIGHT)
add_text_box(slide, Inches(9.0), Inches(2.8), Inches(3.1), Inches(3),
             "• Python 3.10\n\n"
             "• OpenCV (cv2) for image matrices\n\n"
             "• NumPy for fast math\n\n"
             "• Base64 encode for returning images", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 11. UZEM QUESTION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "UZEM Forum Question", "Required component.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(1.5), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(0.5), "Question:", font_size=24, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(2.7), Inches(11), Inches(1),
             "Why is CLAHE (Contrast Limited Adaptive Histogram Equalization) applied to the Luminance (L) channel of the LAB color space rather than applying standard Histogram Equalization to RGB channels?", font_size=18, color=TEXT_LIGHT)


add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(11.33), Inches(2.5), BG_CARD)
add_text_box(slide, Inches(1.2), Inches(4.2), Inches(11), Inches(0.5), "Answer:", font_size=24, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(11), Inches(1.5),
             "Applying equalization directly to R, G, and B changes the color ratios, which alters the original hues (color distortion). By converting to LAB space, we separate Lightness (L) from color data (A, B). Applying CLAHE only to the L channel boosts contrast without changing the actual colors of the fish or water. The 'Limited' part of CLAHE also prevents flat areas like open water from becoming too noisy.", font_size=16, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# 12. CONCLUSION
# -----------------------------------------------------------------------------
slide = new_slide()
slide_header(slide, "Conclusion", "Final thoughts.")

add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(11.33), Inches(4.5), BG_CARD)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.33), Inches(0.5), "Summary of Work:", font_size=24, bold=True, color=ACCENT_OCEAN)
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.33), Inches(3),
             "• Successfully built an underwater image enhancement pipeline using 100% classical Image Processing.\n\n"
             "• Developed an intuitive Node.js/HTML web application featuring Auto-Analysis and Dark Mode.\n\n"
             "• Demonstrated that complex CNNs aren't always necessary for image restoration if mathematical priors (like Dark Channel) are utilized.\n\n"
             "Thanks for listening!", font_size=18, color=TEXT_SEC)

# -----------------------------------------------------------------------------
# SAVING
# -----------------------------------------------------------------------------
def add_border(slide):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height, None, ACCENT_OCEAN)

for i, sl in enumerate(prs.slides):
    if i > 0: # skip title
        # Add a subtle border to all content slides
        pass

out_path = os.path.join(os.path.dirname(__file__), 'AquaClear_Presentation.pptx')
prs.save(out_path)
print(f"Presentation saved to: {out_path}")
print(f"Total slides: {slide_counter[0]}")
