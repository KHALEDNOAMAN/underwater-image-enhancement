"""
AquaClear — SlidesGo Underwater Life Themed Presentation
Dark navy, serif titles, bubble decorations, structured:
Title → Introduction → Main Topic → Results → Conclusion → References
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─── Color Palette (SlidesGo Underwater Life) ─────────────────────────────────
BG_NAVY      = RGBColor(0x00, 0x1C, 0x2E)   # Deep Navy   #001C2E
BG_TEAL      = RGBColor(0x00, 0x3D, 0x5C)   # Mid Teal    #003D5C
ACCENT_CYAN  = RGBColor(0x06, 0xB6, 0xD4)   # Cyan        #06B6D4
ACCENT_OR    = RGBColor(0xF4, 0xA2, 0x61)   # Warm Orange #F4A261
ACCENT_GOLD  = RGBColor(0xFF, 0xD1, 0x66)   # Gold        #FFD166
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT   = RGBColor(0xBF, 0xDB, 0xE8)   # Light Blue-White
TEXT_MUTED   = RGBColor(0x7F, 0xAD, 0xC2)   # Faded blue
CARD_BG      = RGBColor(0x00, 0x2E, 0x46)   # Darker card

# ─── Artifact Paths ───────────────────────────────────────────────────────────
BRAIN_DIR = r"C:\Users\msi\.gemini\antigravity\brain\f59a6a2b-31ec-4512-92ce-db1ebd3f2537"
SAMPLE_DIR = r"c:\Users\msi\OneDrive\Desktop\underwater-image-enhancement\sample_images"
BG_OCEAN  = os.path.join(BRAIN_DIR, "slide_title_background_1774317762365.png")
BG_ACCENT = os.path.join(BRAIN_DIR, "slide_section_accent_1774317780255.png")
IMG_ABSORPTION = os.path.join(BRAIN_DIR, "underwater_light_absorption_chart_1773797577261.png")
IMG_JAFFE      = os.path.join(BRAIN_DIR, "jaffe_mcglamery_model_diagram_1773797596176.png")

# ─── Presentation Setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide_no = [0]

# ─── Helpers ──────────────────────────────────────────────────────────────────
def new_slide(bg_color=BG_NAVY, bg_image=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    if bg_image and os.path.exists(bg_image):
        slide.shapes.add_picture(bg_image, 0, 0, prs.slide_width, prs.slide_height)
        # Semi-transparent dark overlay via XML
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = BG_NAVY
        overlay.line.fill.background()
        # Set alpha to ~78% opacity (val in 1/100 percent: 100000 = fully opaque)
        sp = overlay._element
        solidFill = sp.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha_el.set('val', '78000')
    else:
        fill.solid()
        fill.fore_color.rgb = bg_color
    slide_no[0] += 1
    return slide

def txt(slide, left, top, width, height, text,
        size=20, color=TEXT_WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    r  = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font_name
    p.alignment = align
    return tb

def serif(slide, left, top, width, height, text,
          size=40, color=TEXT_WHITE, bold=True, align=PP_ALIGN.LEFT):
    return txt(slide, left, top, width, height, text,
               size=size, color=color, bold=bold, align=align, font_name="Georgia")

def card(slide, left, top, width, height, color=CARD_BG, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def divider(slide, left, top, width, color=ACCENT_CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def circle_accent(slide, cx, cy, r, color=ACCENT_OR, filled=False):
    left = cx - r
    top  = cy - r
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, r*2, r*2)
    if filled:
        s.fill.solid()
        s.fill.fore_color.rgb = color
    else:
        s.fill.background()
    s.line.color.rgb = color
    s.line.width = Pt(1)
    return s

def num_badge(slide, num_str, left, top):
    """Large decorative section number."""
    serif(slide, left, top, Inches(2), Inches(1.5), num_str,
          size=90, color=ACCENT_CYAN, bold=True)

def slide_num_label(slide):
    txt(slide, prs.slide_width - Inches(0.9), prs.slide_height - Inches(0.45),
        Inches(0.8), Inches(0.4), str(slide_no[0]),
        size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

def section_header_slide(num, title, subtitle=""):
    slide = new_slide(bg_image=BG_ACCENT)
    slide_num_label(slide)
    # Big bubbles decoration
    circle_accent(slide, Inches(11.5), Inches(1.5), Inches(1.8), ACCENT_OR, filled=False)
    circle_accent(slide, Inches(0.8), Inches(6.2), Inches(0.8), ACCENT_CYAN, filled=False)
    circle_accent(slide, Inches(2), Inches(1.2), Inches(0.35), ACCENT_OR, filled=False)
    # Number
    num_badge(slide, num, Inches(0.8), Inches(2))
    divider(slide, Inches(0.8), Inches(3.45), Inches(4), ACCENT_CYAN)
    serif(slide, Inches(0.8), Inches(3.6), Inches(7), Inches(1.2), title,
          size=46, bold=True)
    if subtitle:
        txt(slide, Inches(0.8), Inches(4.8), Inches(7), Inches(0.7), subtitle,
            size=18, color=TEXT_LIGHT)
    return slide

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ① TITLE SLIDE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_OCEAN)
slide_num_label(slide)

# Bubble decorations
circle_accent(slide, Inches(12.2), Inches(1.2), Inches(1.4), ACCENT_OR)
circle_accent(slide, Inches(11.2), Inches(5.8), Inches(0.7), ACCENT_OR)
circle_accent(slide, Inches(1.0), Inches(1.0), Inches(0.5), ACCENT_CYAN)
circle_accent(slide, Inches(0.6), Inches(5.5), Inches(1.2), ACCENT_CYAN)

# Left side colored bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_CYAN; bar.line.fill.background()

# Title block
card(slide, Inches(0.5), Inches(2.2), Inches(9), Inches(3.5), CARD_BG, line=None, radius=True)
serif(slide, Inches(0.8), Inches(2.4), Inches(8.5), Inches(1.5), "AquaClear",
      size=66, bold=True, color=ACCENT_CYAN)
serif(slide, Inches(0.8), Inches(3.8), Inches(8.5), Inches(0.9),
      "Underwater Image Restoration & Enhancement", size=24, bold=False, color=TEXT_LIGHT)
divider(slide, Inches(0.8), Inches(4.75), Inches(5), ACCENT_OR)
txt(slide, Inches(0.8), Inches(4.95), Inches(8.5), Inches(0.5),
    "CENG 455 – Digital Image Processing  ·  Istanbul Arel University", size=14, color=TEXT_MUTED)
txt(slide, Inches(0.8), Inches(5.45), Inches(8.5), Inches(0.5),
    "Khaled Noaman  ·  200303812  ·  Spring 2026", size=14, color=TEXT_WHITE, bold=True)

# Show coral reef sample image in circle-ish on the right
coral = os.path.join(SAMPLE_DIR, 'coral_reef.jpg')
if os.path.exists(coral):
    slide.shapes.add_picture(coral, Inches(9.8), Inches(1.0), width=Inches(3.1), height=Inches(5.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ② TABLE OF CONTENTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12), Inches(0.9), Inches(1.1), ACCENT_GOLD)
circle_accent(slide, Inches(1.2), Inches(6.5), Inches(0.6), ACCENT_CYAN)

serif(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.8), "Contents", size=38, bold=True)
divider(slide, Inches(0.8), Inches(1.25), Inches(7), ACCENT_CYAN)

sections = [
    ("01", "Introduction",  "Physics of underwater light degradation"),
    ("02", "Main Topic",    "The AquaClear Processing Pipeline"),
    ("03", "Results",       "Visual Before & After Comparisons"),
    ("04", "Conclusion",    "Key Takeaways and Future Work"),
    ("05", "References",    "Academic Sources"),
]
for i, (num, title, sub) in enumerate(sections):
    top = Inches(1.5) + i * Inches(1.04)
    serif(slide, Inches(0.8), top, Inches(0.8), Inches(0.7), num, size=28, color=ACCENT_CYAN, bold=True)
    txt(slide, Inches(1.9), top, Inches(5.5), Inches(0.5), title, size=22, color=TEXT_WHITE, bold=True)
    txt(slide, Inches(1.9), top + Inches(0.38), Inches(7), Inches(0.45), sub, size=14, color=TEXT_MUTED)
    divider(slide, Inches(1.7), top + Inches(0.9), Inches(8.5), TEXT_MUTED)

# Image stack on the right
diver = os.path.join(SAMPLE_DIR, 'diver.jpg')
if os.path.exists(diver):
    slide.shapes.add_picture(diver, Inches(10.2), Inches(0.5), width=Inches(2.7), height=Inches(6.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ③ SECTION 01 HEADER — INTRODUCTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_header_slide("01", "Introduction", "The physics of underwater light degradation")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ④ Why Underwater Images Are Degraded
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12.8), Inches(6.8), Inches(1.5), ACCENT_OR)
circle_accent(slide, Inches(0.5), Inches(0.5), Inches(0.4), ACCENT_CYAN)

serif(slide, Inches(0.7), Inches(0.3), Inches(8), Inches(0.7), "Why Underwater Images Are Degraded", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(6.5), ACCENT_CYAN)

facts = [
    (ACCENT_OR,   "Color Attenuation",  "Red light (650nm) is completely absorbed within just 3–5m of depth. Only blue/green light survives, creating the pervasive cyan cast seen in all raw underwater photos."),
    (ACCENT_CYAN, "Backscattering",     "Ambient light reflects off suspended particles (marine snow, plankton) directly into the camera lens, creating a dense gray-blue fog that eliminates visibility."),
    (ACCENT_GOLD, "Forward Scattering", "Light reflecting off the subject is scattered by particles before reaching the lens, blurring edges and destroying high-frequency details."),
]
for i, (color, title, body) in enumerate(facts):
    top = Inches(1.3) + i * Inches(1.7)
    card(slide, Inches(0.7), top, Inches(6.5), Inches(1.55), CARD_BG, line=color)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), top + Inches(0.5), Inches(0.35), Inches(0.35))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(slide, Inches(1.3), top + Inches(0.15), Inches(6), Inches(0.45), title, size=18, color=color, bold=True)
    txt(slide, Inches(1.3), top + Inches(0.6), Inches(6), Inches(0.8), body, size=13, color=TEXT_LIGHT)

# Right side: light absorption chart
if os.path.exists(IMG_ABSORPTION):
    slide.shapes.add_picture(IMG_ABSORPTION, Inches(7.6), Inches(1.1), width=Inches(5.4))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑤ Jaffe-McGlamery Optical Model
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "The Jaffe-McGlamery Optical Model", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(7), ACCENT_CYAN)

txt(slide, Inches(0.7), Inches(1.25), Inches(5.8), Inches(1.1),
    "The dominant model for underwater image formation states that the radiance captured by a camera sensor is the sum of three physical components:", size=15, color=TEXT_LIGHT)
txt(slide, Inches(0.7), Inches(2.4), Inches(5.8), Inches(0.5),
    "E_total = E_direct + E_forward_scatter + E_backscatter", size=16, color=ACCENT_GOLD, bold=True)

components = [
    ("Direct Signal",      "Unscattered light from the object — the 'true' signal we want.",    ACCENT_CYAN),
    ("Forward Scatter",    "Deflected signal rays causing blur at the subject level.",            ACCENT_OR),
    ("Backscatter (Haze)", "Ambient particles deflecting light directly into the lens as fog.",   ACCENT_GOLD),
]
for i, (title, body, color) in enumerate(components):
    top = Inches(3.0) + i * Inches(1.2)
    divider(slide, Inches(0.7), top, Inches(5.8), color)
    txt(slide, Inches(0.7), top + Inches(0.07), Inches(5.8), Inches(0.4), title, size=17, color=color, bold=True)
    txt(slide, Inches(0.7), top + Inches(0.5), Inches(5.8), Inches(0.55), body, size=13, color=TEXT_LIGHT)

if os.path.exists(IMG_JAFFE):
    slide.shapes.add_picture(IMG_JAFFE, Inches(7.1), Inches(1.1), width=Inches(5.9))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑥ Classical vs Deep Learning
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12.5), Inches(0.8), Inches(1.0), ACCENT_GOLD)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Classical IP vs Deep Learning", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(7), ACCENT_CYAN)

# Two column comparison
headers = ["Classical Image Processing ✓", "Deep Learning / CNNs ✗"]
cols_left = [Inches(0.7), Inches(6.8)]
for ci, (hdr, left) in enumerate(zip(headers, cols_left)):
    col = ACCENT_CYAN if ci == 0 else ACCENT_OR
    card(slide, left, Inches(1.25), Inches(5.7), Inches(5.5), CARD_BG, line=col)
    txt(slide, left + Inches(0.2), Inches(1.4), Inches(5.3), Inches(0.55), hdr, size=20, color=col, bold=True)

rows = [
    ("✓ Zero hallucinations — no invented pixels",           "✗ Frequently hallucinates textures"),
    ("✓ Needs no training data",                              "✗ Requires huge paired raw/clear datasets"),
    ("✓ Fully transparent and mathematically auditable",      "✗ Black-box, unexplainable outputs"),
    ("✓ < 0.5 sec on CPU, no GPU required",                  "✗ Requires GPU, minutes per frame"),
    ("✓ Safe for scientific measurement",                     "✗ Risky for biological classification"),
]
for i, (left_text, right_text) in enumerate(rows):
    top = Inches(2.1) + i * Inches(0.85)
    txt(slide, Inches(0.9), top, Inches(5.3), Inches(0.75), left_text, size=14, color=TEXT_LIGHT)
    txt(slide, Inches(7.0), top, Inches(5.3), Inches(0.75), right_text, size=14, color=TEXT_MUTED)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑦ SECTION 02 HEADER — MAIN TOPIC
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_header_slide("02", "Main Topic", "The AquaClear Processing Pipeline — Step by Step")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑧ Pipeline Overview
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12.5), Inches(6.2), Inches(1.3), ACCENT_OR)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "The AquaClear Pipeline", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(7), ACCENT_CYAN)

steps = [
    ("1", "Auto-Analysis",     "Scans the image matrix to detect blur, haze, and color deficit levels, then sets optimal parameters.", ACCENT_CYAN),
    ("2", "Noise Reduction",   "Non-linear Bilateral Filter selectively smooths flat areas while protecting sharp edges.",              ACCENT_OR),
    ("3", "Color Correction",  "Gray World Theorem scales the Red channel gain to neutralize the blue-green cast.",                    ACCENT_GOLD),
    ("4", "Dehazing",          "Dark Channel Prior estimates & subtracts backscattered transmission fog.",                             ACCENT_CYAN),
    ("5", "Contrast (CLAHE)",  "Adapative Histogram Equalization on the L channel of LAB color space.",                               ACCENT_OR),
    ("6", "Sharpening",        "Unsharp Masking boosts high-frequency edge data to restore detail lost to forward scattering.",        ACCENT_GOLD),
]
for i, (num, title, body, color) in enumerate(steps):
    col = 0 if i < 3 else 1
    row = i % 3
    left = Inches(0.7) + col * Inches(6.3)
    top  = Inches(1.3) + row * Inches(1.9)
    card(slide, left, top, Inches(5.9), Inches(1.75), CARD_BG, line=color)
    serif(slide, left + Inches(0.2), top + Inches(0.1), Inches(0.5), Inches(0.6), num, size=32, color=color, bold=True)
    txt(slide, left + Inches(0.8), top + Inches(0.1), Inches(5), Inches(0.5), title, size=18, color=color, bold=True)
    txt(slide, left + Inches(0.8), top + Inches(0.65), Inches(4.9), Inches(0.9), body, size=13, color=TEXT_LIGHT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑨ Auto-Analysis Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Phase 1: Auto-Analysis Heuristics", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(6.5), ACCENT_CYAN)
txt(slide, Inches(0.7), Inches(1.2), Inches(12), Inches(0.6),
    "Static pipelines fail because every ocean differs. AquaClear dynamically calibrates for each image:", size=16, color=TEXT_LIGHT)

metrics = [
    ("Red Deficit",        "b_mean - r_mean → if > 30: aggressive color correction",            ACCENT_OR),
    ("Haze Level",         "Dark Channel Mean / 255 → if > 0.4: omega = 0.98 (dense fog)",      ACCENT_CYAN),
    ("Blur Level",         "Laplacian Variance → if < 100: switch to NL-Means + heavy sharpening", ACCENT_GOLD),
    ("Contrast",           "Grayscale Std Dev → if < 20: CLAHE Clip = 5.0 (very low contrast)", ACCENT_CYAN),
]
for i, (metric, logic, color) in enumerate(metrics):
    top = Inches(2.0) + i * Inches(1.2)
    col = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), top, Inches(0.08), Inches(0.9))
    col.fill.solid(); col.fill.fore_color.rgb = color; col.line.fill.background()
    txt(slide, Inches(0.95), top, Inches(5), Inches(0.45), metric, size=18, color=color, bold=True)
    txt(slide, Inches(0.95), top + Inches(0.47), Inches(11), Inches(0.55), logic, size=14, color=TEXT_LIGHT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑩ Color Correction Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Phase 3: Color Correction", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(5), ACCENT_CYAN)
txt(slide, Inches(0.7), Inches(1.2), Inches(6), Inches(0.6),
    "Gray World Theorem: In a scene with diversity, average R = G = B = Gray", size=16, color=TEXT_LIGHT)

card(slide, Inches(0.7), Inches(1.9), Inches(6), Inches(2.5), CARD_BG, line=ACCENT_OR)
txt(slide, Inches(0.9), Inches(2.1), Inches(5.8), Inches(0.5), "Algorithm:", size=18, color=ACCENT_OR, bold=True)
formula_lines = [
    "r_mean, g_mean, b_mean = mean of each channel",
    "total_mean = (r_mean + g_mean + b_mean) / 3",
    "r_gain = total_mean / r_mean  ← (boosts red)",
    "R_corrected = clip(R * r_gain, 0, 255)",
]
for i, line in enumerate(formula_lines):
    txt(slide, Inches(0.9), Inches(2.65) + i*Inches(0.42), Inches(5.8), Inches(0.42),
        line, size=14, color=ACCENT_GOLD)

txt(slide, Inches(0.7), Inches(4.55), Inches(6), Inches(1.8),
    "This single operation removes the dominant cyan tint by digitally restoring the wavelengths that water physically absorbed. Without it, the dehazing step cannot function correctly.",
    size=14, color=TEXT_LIGHT)

reef = os.path.join(SAMPLE_DIR, 'tropical_fish.jpg')
if os.path.exists(reef):
    slide.shapes.add_picture(reef, Inches(7.2), Inches(1.1), width=Inches(5.7))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑪ Dehazing Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Phase 4: Dark Channel Prior Dehazing", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(6), ACCENT_CYAN)

dcp_steps = [
    ("Step 1", "Compute Dark Channel",    "D(x) = min of pixel across all 3 channels in a local patch",   ACCENT_CYAN),
    ("Step 2", "Estimate Atmospheric A",  "Brightest 0.1% of pixels in the Dark Channel = Airlight (A)",  ACCENT_OR),
    ("Step 3", "Calculate Transmission",  "t(x) = 1 – ω × min(I(y)/A)  |  ω controls strength",          ACCENT_GOLD),
    ("Step 4", "Recover Scene Radiance",  "J(x) = (I(x) – A) / max(t(x), 0.1) + A",                      ACCENT_CYAN),
]
for i, (step, title, formula, color) in enumerate(dcp_steps):
    top = Inches(1.4) + i * Inches(1.4)
    card(slide, Inches(0.7), top, Inches(7.2), Inches(1.25), CARD_BG, line=color)
    txt(slide, Inches(0.9), top + Inches(0.1), Inches(1.2), Inches(0.45), step, size=14, color=color, bold=True)
    txt(slide, Inches(2.1), top + Inches(0.1), Inches(5.5), Inches(0.45), title, size=17, color=TEXT_WHITE, bold=True)
    txt(slide, Inches(0.9), top + Inches(0.62), Inches(6.8), Inches(0.5), formula, size=13, color=ACCENT_GOLD)

deep = os.path.join(SAMPLE_DIR, 'deep_ocean.jpg')
if os.path.exists(deep):
    slide.shapes.add_picture(deep, Inches(8.5), Inches(1.1), width=Inches(4.5))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑫ CLAHE & LAB Deep Dive
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Phase 5: CLAHE on LAB L-Channel", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(5.5), ACCENT_CYAN)

card(slide, Inches(0.7), Inches(1.25), Inches(5.9), Inches(5.5), CARD_BG, line=ACCENT_GOLD)
txt(slide, Inches(0.9), Inches(1.4), Inches(5.6), Inches(0.5), "Why LAB color space?", size=19, color=ACCENT_GOLD, bold=True)
txt(slide, Inches(0.9), Inches(1.95), Inches(5.6), Inches(1.5),
    "LAB separates Luminance (L) from color (A, B channels).\nEquilizing R, G, B independently shifts color ratios and creates false hues. Equalizing only L preserves the restored colors from Phase 3.",
    size=14, color=TEXT_LIGHT)
divider(slide, Inches(0.9), Inches(3.5), Inches(5.5), TEXT_MUTED)
txt(slide, Inches(0.9), Inches(3.6), Inches(5.6), Inches(0.5), "Why CLAHE not Global HE?", size=19, color=ACCENT_CYAN, bold=True)
txt(slide, Inches(0.9), Inches(4.15), Inches(5.6), Inches(1.8),
    "Global HE stretches the entire image uniformly — in dark underwater scenes it creates massive noise artifacts in flat water regions.\n\nCLAHE divides the image into 8×8 tiles and limits the histogram peak (Clip Limit), boosting local contrast while preventing noise explosion.",
    size=14, color=TEXT_LIGHT)

jelly = os.path.join(SAMPLE_DIR, 'jellyfish.jpg')
if os.path.exists(jelly):
    slide.shapes.add_picture(jelly, Inches(7.2), Inches(1.1), width=Inches(5.7))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑬ System Architecture
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12.8), Inches(7.0), Inches(1.5), ACCENT_OR)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "System Architecture", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(6), ACCENT_CYAN)

arch = [
    ("HTML/CSS/JS\nFrontend", "Upload UI, Dark Mode, Before/After sliders, Metric gauges", ACCENT_CYAN),
    ("Node.js\nExpress Server", "Multi-part upload, child-process forking, JSON routing", ACCENT_OR),
    ("Python\nIP Engine", "NumPy / OpenCV matrix processing. All 6 pipeline phases.", ACCENT_GOLD),
]
for i, (title, body, color) in enumerate(arch):
    left = Inches(0.7) + i * Inches(4.15)
    card(slide, left, Inches(1.4), Inches(3.8), Inches(4.5), CARD_BG, line=color)
    serif(slide, left + Inches(0.2), Inches(1.6), Inches(3.5), Inches(1.0), title, size=22, color=color, bold=True)
    divider(slide, left + Inches(0.2), Inches(2.75), Inches(3.3), color)
    txt(slide, left + Inches(0.2), Inches(2.9), Inches(3.5), Inches(2.5), body, size=15, color=TEXT_LIGHT)
    if i < len(arch)-1:
        txt(slide, left + Inches(3.9), Inches(3.3), Inches(0.4), Inches(0.6), "→", size=30, color=TEXT_MUTED)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑭ SECTION 03 HEADER — RESULTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_header_slide("03", "Results", "Visual Before & After Restoration Comparisons")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑮–⑲ Before/After Image Slides
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
ba_pairs = [
    ("Coral Reef",       "coral_reef.jpg",      "coral_reef_enhanced.jpg"),
    ("Jellyfish",        "jellyfish.jpg",        "jellyfish_enhanced.jpg"),
    ("Deep Ocean Wreck", "deep_ocean.jpg",       "deep_ocean_enhanced.jpg"),
    ("Tropical Fish",    "tropical_fish.jpg",    "tropical_fish_enhanced.jpg"),
    ("Underwater Reef",  "underwater_reef.jpg",  "underwater_reef_enhanced.jpg"),
]
for label, raw_fn, enh_fn in ba_pairs:
    slide = new_slide(bg_image=BG_ACCENT)
    slide_num_label(slide)
    circle_accent(slide, Inches(12.9), Inches(0.7), Inches(0.9), ACCENT_OR)

    serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.65),
          f"Results: {label}", size=28, bold=True)
    divider(slide, Inches(0.7), Inches(1.0), Inches(12), ACCENT_CYAN)

    raw_path = os.path.join(SAMPLE_DIR, raw_fn)
    enh_path = os.path.join(SAMPLE_DIR, enh_fn)
    
    img_h = Inches(4.8)
    img_w = Inches(5.8)
    if os.path.exists(raw_path):
        slide.shapes.add_picture(raw_path, Inches(0.5), Inches(1.5), width=img_w, height=img_h)
    card(slide, Inches(0.5), Inches(6.35), Inches(5.8), Inches(0.65), BG_TEAL, line=ACCENT_OR)
    txt(slide, Inches(0.7), Inches(6.45), Inches(5.5), Inches(0.5), "BEFORE  –  Raw Underwater Input", size=14, color=ACCENT_OR, bold=True)

    if os.path.exists(enh_path):
        slide.shapes.add_picture(enh_path, Inches(6.8), Inches(1.5), width=img_w, height=img_h)
    card(slide, Inches(6.8), Inches(6.35), Inches(5.8), Inches(0.65), BG_TEAL, line=ACCENT_CYAN)
    txt(slide, Inches(7.0), Inches(6.45), Inches(5.5), Inches(0.5), "AFTER  –  AquaClear Enhanced", size=14, color=ACCENT_CYAN, bold=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ⑳ Objective Metrics
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12.5), Inches(6.8), Inches(1.3), ACCENT_GOLD)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Objective Quality Evaluation Metrics", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(7), ACCENT_CYAN)

metrics_data = [
    ("PSNR",        "Peak Signal-to-Noise Ratio",     "10·log₁₀(MAX² / MSE)", "20–35 dB  (Higher = Better Fidelity)", ACCENT_CYAN),
    ("SSIM",        "Structural Similarity Index",     "Luminance × Contrast × Structure",  "0.80–0.95  (Higher = Less Distortion)", ACCENT_OR),
    ("Entropy",     "Shannon Information Density",     "–Σ P(i)·log₂P(i)",     "+15–25% increase (More scene detail)", ACCENT_GOLD),
    ("Colorfulness","C = √(σRG² + σYB²) + 0.3·√(μRG² + μYB²)", "Per-channel saturation variance", "Strong increase vs raw cyan image", ACCENT_CYAN),
]
col_w = Inches(5.7)
for i, (name, full, formula, result, color) in enumerate(metrics_data):
    col = i % 2
    row = i // 2
    left = Inches(0.7) + col * Inches(6.3)
    top  = Inches(1.35) + row * Inches(2.7)
    card(slide, left, top, col_w, Inches(2.5), CARD_BG, line=color)
    serif(slide, left+Inches(0.2), top+Inches(0.1), Inches(1.5), Inches(0.6), name, size=28, color=color, bold=True)
    txt(slide, left+Inches(0.2), top+Inches(0.7), col_w-Inches(0.4), Inches(0.5), full, size=14, color=TEXT_LIGHT, bold=True)
    txt(slide, left+Inches(0.2), top+Inches(1.15), col_w-Inches(0.4), Inches(0.55), formula, size=12, color=ACCENT_GOLD)
    divider(slide, left+Inches(0.2), top+Inches(1.75), col_w-Inches(0.4), color)
    txt(slide, left+Inches(0.2), top+Inches(1.87), col_w-Inches(0.4), Inches(0.5), result, size=13, color=TEXT_MUTED)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ㉑ SECTION 04 HEADER — CONCLUSION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_header_slide("04", "Conclusion", "Achievements, Limitations, and Future Directions")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ㉒ Conclusion Slide
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "Conclusion", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(6), ACCENT_CYAN)

# Key achievements
achieves = [
    "Successfully engineered a full 6-phase classical image processing pipeline tailored to aquatic optics.",
    "Achieved real-time restoration (< 0.5 seconds) on standard CPU hardware — no GPU required.",
    "Auto-analysis heuristic engine adaptively tunes all parameters per image, removing the need for manual tuning.",
    "Objective metrics (PSNR, SSIM, Entropy) confirmed structural fidelity and genuine information recovery.",
    "Wrapped the engine in a modern dark-themed web application with live Before/After visualization.",
]
for i, line in enumerate(achieves):
    top = Inches(1.25) + i * Inches(0.98)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.15), Inches(0.22), Inches(0.22))
    s.fill.solid(); s.fill.fore_color.rgb = ACCENT_CYAN; s.line.fill.background()
    txt(slide, Inches(1.1), top, Inches(6.5), Inches(0.85), line, size=14, color=TEXT_LIGHT)

# Limitations & Future
card(slide, Inches(7.9), Inches(1.2), Inches(5.1), Inches(5.6), CARD_BG, line=ACCENT_OR)
txt(slide, Inches(8.1), Inches(1.35), Inches(4.8), Inches(0.5), "Limitations & Future Work", size=17, color=ACCENT_OR, bold=True)
divider(slide, Inches(8.1), Inches(1.9), Inches(4.6), ACCENT_OR)
limitations = [
    "Gray World fails on monochromatic single-subject shots.",
    "DCP struggles with artificial dive lights.",
    "",
    "Future: Stereo depth mapping for per-pixel attenuation.",
    "Future: CUDA port for live 60fps video telemetry.",
    "Future: Mobile deployment for handheld dive cams.",
]
for i, line in enumerate(limitations):
    color = ACCENT_OR if line.startswith("Gray") or line.startswith("DCP") else (ACCENT_CYAN if line.startswith("Future") else TEXT_MUTED)
    txt(slide, Inches(8.1), Inches(2.0) + i*Inches(0.78), Inches(4.8), Inches(0.7), line, size=13, color=color)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ㉓ SECTION 05 HEADER — REFERENCES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
section_header_slide("05", "References", "Academic & Technical Sources")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ㉔ References Slide
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_ACCENT)
slide_num_label(slide)
circle_accent(slide, Inches(12.5), Inches(6.5), Inches(1.2), ACCENT_GOLD)

serif(slide, Inches(0.7), Inches(0.3), Inches(10), Inches(0.7), "References", size=30, bold=True)
divider(slide, Inches(0.7), Inches(1.05), Inches(10), ACCENT_CYAN)

refs = [
    ("[1]", "R. C. Gonzalez and R. E. Woods", "Digital Image Processing, 4th ed.", "Pearson, 2018"),
    ("[2]", "G. Bradski", "The OpenCV Library", "Dr. Dobb's Journal of Software Tools, 2000"),
    ("[3]", "K. He, J. Sun, and X. Tang", "Single Image Haze Removal Using Dark Channel Prior",  "IEEE Transactions on PAMI, 2011"),
    ("[4]", "J. S. Jaffe", "Computer Modeling and Design of Optimal Underwater Imaging Systems", "IEEE Journal of Oceanic Engineering, 1990"),
    ("[5]", "A. M. Reza", "Realization of CLAHE for Real-Time Image Enhancement", "Journal of VLSI Signal Processing, 2004"),
]
for i, (num, author, title, pub) in enumerate(refs):
    top = Inches(1.4) + i * Inches(1.1)
    divider(slide, Inches(0.7), top, Inches(12), BG_TEAL)
    txt(slide, Inches(0.7), top + Inches(0.08), Inches(0.5), Inches(0.6), num, size=13, color=ACCENT_CYAN, bold=True)
    txt(slide, Inches(1.3), top + Inches(0.08), Inches(11.5), Inches(0.4), author + " — " + title, size=15, color=TEXT_WHITE, bold=True)
    txt(slide, Inches(1.3), top + Inches(0.55), Inches(11.5), Inches(0.45), pub, size=13, color=TEXT_MUTED)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ㉕ FINAL: Thank You
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = new_slide(bg_image=BG_OCEAN)
slide_num_label(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_CYAN; bar.line.fill.background()

circle_accent(slide, Inches(11.5), Inches(1.0), Inches(1.5), ACCENT_OR)
circle_accent(slide, Inches(10.5), Inches(6.5), Inches(0.7), ACCENT_CYAN)
circle_accent(slide, Inches(1.5), Inches(6.3), Inches(0.5), ACCENT_GOLD)

card(slide, Inches(0.5), Inches(2.5), Inches(8.5), Inches(3.5), CARD_BG, radius=True)
serif(slide, Inches(1), Inches(2.7), Inches(7.5), Inches(1.2), "Thank You", size=64, color=ACCENT_CYAN, bold=True)
divider(slide, Inches(1), Inches(3.85), Inches(5), ACCENT_OR)
txt(slide, Inches(1), Inches(4.05), Inches(7.5), Inches(0.55), "Questions & Discussion", size=24, color=TEXT_LIGHT)
txt(slide, Inches(1), Inches(4.70), Inches(7.5), Inches(0.5), "Khaled Noaman  •  200303812", size=16, color=TEXT_MUTED)
txt(slide, Inches(1), Inches(5.10), Inches(7.5), Inches(0.5), "CENG 455  •  Istanbul Arel University  •  Spring 2026", size=14, color=TEXT_MUTED)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out = os.path.join(r"c:\Users\msi\OneDrive\Desktop\underwater-image-enhancement",
                   "AquaClear_SlidesGo_Theme.pptx")
prs.save(out)
print(f"Saved: {out}\nTotal slides: {slide_no[0]}")
